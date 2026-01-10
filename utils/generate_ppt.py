import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation(md_file_path, output_pptx_path):
    prs = Presentation()

    # Read the markdown content
    with open(md_file_path, 'r', encoding='utf-8') as f:
        content = f.read()

    # 1. Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Crypto Trading AI: Project Flow"
    subtitle.text = "From Data Collection to Automation"

    # 2. Flowchart (Text Representation)
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = "Project Lifecycle Flow"
    tf = body_shape.text_frame
    tf.text = "1. Data Collection"
    
    p = tf.add_paragraph()
    p.text = "Fetch OHLCV data -> Raw Market Data"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "2. Feature Engineering"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Calculate Indicators -> Engineered Features"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "3. Prediction & ML"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Train Models (XGBoost/LightGBM) -> Predictions"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "4. Automation"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Control Center triggers Collection & Prediction -> Signals/Alerts"
    p.level = 1

    # 3. Detailed Breakdown
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Key Components & Files"

    tf = body_shape.text_frame
    tf.text = "Data Collection"
    p = tf.add_paragraph()
    p.text = "File: comprehensive_ml_collector_v2.py"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Fetches historical data into SQLite/CSV."
    p.level = 2

    p = tf.add_paragraph()
    p.text = "Feature Engineering"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "File: crypto_ai/features/engineer.py"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Calculates technical indicators for ML models."
    p.level = 2

    p = tf.add_paragraph()
    p.text = "Prediction System"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "File: optimized_ml_system.py"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Automation Control"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "File: crypto_control_center.py"
    p.level = 1

    # 4. Mermaid Diagram Code (Reference)
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Mermaid Diagram Code"
    
    # Extract mermaid code block
    try:
        start = content.find("```mermaid")
        end = content.find("```", start + 10)
        if start != -1 and end != -1:
            mermaid_code = content[start+10:end].strip()
            
            # Add text box for code
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(9)
            height = Inches(5.5)
            
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            
            p = tf.add_paragraph()
            p.text = mermaid_code
            p.font.size = Pt(10)
            p.font.name = "Courier New"
            
            # Remove the default placeholder text if we added a custom textbox
            # or just leave it blank
    except:
        pass

    # Save
    prs.save(output_pptx_path)
    print(f"Presentation saved to {output_pptx_path}")

if __name__ == "__main__":
    md_file = "/home/ofri/.gemini/antigravity/brain/1d993db8-b45d-4918-b1c4-d563e9887588/project_flowchart.md"
    output_file = "/home/ofri/.gemini/antigravity/brain/1d993db8-b45d-4918-b1c4-d563e9887588/Project_Flowchart.pptx"
    create_presentation(md_file, output_file)
